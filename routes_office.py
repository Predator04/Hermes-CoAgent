"""Office suite authoring routes (Excel / Word / PowerPoint).

File engine (headless, works without Office installed):
  POST /office/excel/create - create a new .xlsx workbook from {path, sheets:[{name, rows:[...]}]}
  POST /office/excel/set    - set cell values in an existing .xlsx and save
  GET  /office/excel/read   - read a sheet (or range) back as a grid
  POST /office/word/create  - create a .docx from {path, title, paragraphs:[...]}
  GET  /office/word/read    - read paragraphs from a .docx
  POST /office/ppt/create   - create a .pptx from {path, slides:[{title, bullets:[...]}]}
  GET  /office/ppt/read     - read slide titles + body text from a .pptx

COM engine (live desktop Office, Windows only):
  POST /office/excel/com/open - open a workbook in the running Excel application
  POST /office/excel/com/set  - set a cell value in the live workbook (and save)
  POST /office/excel/com/close- close a live workbook / quit Excel

All heavy deps (openpyxl, python-docx, python-pptx, win32com) are imported
inside try/except so the Linux syntax-check CI stays green and a missing
optional module returns a clean 501 with an install hint instead of a crash.
"""

import json
import os
from pathlib import Path

from flask import jsonify

from shared import _json_body, _log, _missing_field

# ---------------------------------------------------------------------------
# Optional dependency probes. Each is guarded so a missing module never breaks
# import; endpoints return 501 with a clear hint when their engine is absent.
# ---------------------------------------------------------------------------
try:
    import openpyxl  # type: ignore
    _HAS_XLSX = True
except Exception:
    openpyxl = None
    _HAS_XLSX = False

try:
    import docx  # type: ignore
    _HAS_DOCX = True
except Exception:
    docx = None
    _HAS_DOCX = False

try:
    import pptx  # type: ignore
    from pptx.util import Inches  # type: ignore
    _HAS_PPTX = True
except Exception:
    pptx = None
    Inches = None
    _HAS_PPTX = False

try:
    import win32com.client  # type: ignore
    _HAS_COM = True
except Exception:
    win32com = None
    _HAS_COM = False


def _missing_engine(engine, pkg):
    return jsonify({
        "error": f"{engine} engine unavailable",
        "hint": f"Install the optional dependency: pip install {pkg}",
    }), 501


def _resolve_path(raw):
    if not isinstance(raw, str) or not raw.strip():
        return None
    try:
        return str(Path(os.path.expanduser(raw)).resolve())
    except (OSError, ValueError):
        return None


# ---------------------------------------------------------------------------
# Excel (file engine)
# ---------------------------------------------------------------------------
def _excel_create(path, sheets):
    if not _HAS_XLSX:
        return None, _missing_engine("Excel", "openpyxl")
    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # drop the default sheet so sheet order matches input
    if not sheets:
        wb.create_sheet("Sheet1")
    else:
        for sheet in sheets:
            if not isinstance(sheet, dict):
                continue
            name = str(sheet.get("name") or "Sheet").strip() or "Sheet"
            ws = wb.create_sheet(title=name[:31])
            rows = sheet.get("rows") or sheet.get("data") or []
            for r_idx, row in enumerate(rows, start=1):
                if isinstance(row, (list, tuple)):
                    for c_idx, value in enumerate(row, start=1):
                        ws.cell(row=r_idx, column=c_idx, value=value)
                else:
                    ws.cell(row=r_idx, column=1, value=row)
    wb.save(path)
    return {"path": path, "sheets": wb.sheetnames}, None


def _excel_set(path, updates):
    if not _HAS_XLSX:
        return None, _missing_engine("Excel", "openpyxl")
    wb = openpyxl.load_workbook(path)
    changed = 0
    # updates is either {"Sheet1": {"A1": v, ...}} or a flat {cell: value} with
    # an optional "sheet" key.
    sheet_name = str(updates.get("sheet") or "Sheet1") if isinstance(updates, dict) else "Sheet1"
    flat = updates.get("cells") if isinstance(updates, dict) and "cells" in updates else updates
    if isinstance(flat, dict):
        for cell_ref, value in flat.items():
            if cell_ref in ("sheet", "cells"):
                continue
            if isinstance(value, dict):
                # per-sheet style: {"Sheet2": {"B2": v}}
                for sub_sheet, sub_cells in value.items():
                    if not isinstance(sub_cells, dict):
                        continue
                    ws = wb[sub_sheet] if sub_sheet in wb.sheetnames else wb.create_sheet(sub_sheet)
                    for ref, v in sub_cells.items():
                        ws[ref] = v
                        changed += 1
                continue
            ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.create_sheet(sheet_name)
            ws[cell_ref] = value
            changed += 1
    wb.save(path)
    return {"path": path, "changed_cells": changed}, None


def _excel_read(path, sheet=None, cell_range=None):
    if not _HAS_XLSX:
        return None, _missing_engine("Excel", "openpyxl")
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
    if cell_range:
        cells = ws[cell_range]
        grid = [[c.value for c in row] for row in cells]
    else:
        grid = [[c.value for c in row] for row in ws.iter_rows()]
    return {"path": path, "sheet": ws.title, "grid": grid, "rows": len(grid)}, None


# ---------------------------------------------------------------------------
# Word (file engine)
# ---------------------------------------------------------------------------
def _word_create(path, title, paragraphs):
    if not _HAS_DOCX:
        return None, _missing_engine("Word", "python-docx")
    document = docx.Document()
    if title:
        document.add_heading(str(title), level=0)
    for para in paragraphs or []:
        document.add_paragraph(str(para))
    document.save(path)
    return {"path": path, "paragraphs": len(paragraphs or []) + (1 if title else 0)}, None


def _word_read(path):
    if not _HAS_DOCX:
        return None, _missing_engine("Word", "python-docx")
    document = docx.Document(path)
    paragraphs = [p.text for p in document.paragraphs]
    return {"path": path, "paragraphs": paragraphs, "count": len(paragraphs)}, None


# ---------------------------------------------------------------------------
# PowerPoint (file engine)
# ---------------------------------------------------------------------------
def _ppt_create(path, slides):
    if not _HAS_PPTX:
        return None, _missing_engine("PowerPoint", "python-pptx")
    prs = pptx.Presentation()
    for slide in slides or []:
        if not isinstance(slide, dict):
            continue
        layout = prs.slide_layouts[1]  # Title and Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = str(slide.get("title") or "")
        body = slide.get("bullets") or slide.get("body") or []
        if isinstance(body, str):
            body = [body]
        body_frame = s.placeholders[1]
        tf = body_frame.text_frame
        tf.clear()
        for i, bullet in enumerate(body):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(bullet)
    prs.save(path)
    return {"path": path, "slides": len(slides or [])}, None


def _ppt_read(path):
    if not _HAS_PPTX:
        return None, _missing_engine("PowerPoint", "python-pptx")
    prs = pptx.Presentation(path)
    slides = []
    for s in prs.slides:
        title = s.shapes.title.text if s.shapes.title else ""
        bullets = []
        for shape in s.shapes:
            if shape.has_text_frame and shape != s.shapes.title:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        bullets.append(para.text)
        slides.append({"title": title, "bullets": bullets})
    return {"path": path, "slides": slides, "count": len(slides)}, None


def register_routes(app, state, require_auth):
    # ------------------------- Excel file engine -------------------------
    @app.route("/office/excel/create", methods=["POST"])
    @require_auth
    def route_office_excel_create():
        body = _json_body()
        if not isinstance(body, dict):
            body = {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        sheets = body.get("sheets") or body.get("data")
        result, err = _excel_create(path, sheets)
        if err:
            return err
        _log(f"office/excel/create path={path}")
        return jsonify({"status": "ok", **result})

    @app.route("/office/excel/set", methods=["POST"])
    @require_auth
    def route_office_excel_set():
        body = _json_body()
        if not isinstance(body, dict):
            body = {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        if not os.path.isfile(path):
            return jsonify({"error": "file not found", "path": path}), 404
        updates = body.get("updates") or body.get("cells") or {}
        result, err = _excel_set(path, dict(updates, sheet=body.get("sheet", "Sheet1")) if isinstance(updates, dict) else updates)
        if err:
            return err
        _log(f"office/excel/set path={path}")
        return jsonify({"status": "ok", **result})

    @app.route("/office/excel/read", methods=["GET", "POST"])
    @require_auth
    def route_office_excel_read():
        body = _json_body()
        if not isinstance(body, dict):
            body = {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        if not os.path.isfile(path):
            return jsonify({"error": "file not found", "path": path}), 404
        result, err = _excel_read(path, body.get("sheet"), body.get("range"))
        if err:
            return err
        return jsonify(result)

    # ------------------------- Word file engine --------------------------
    @app.route("/office/word/create", methods=["POST"])
    @require_auth
    def route_office_word_create():
        body = _json_body()
        if not isinstance(body, dict):
            body = {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        result, err = _word_create(path, body.get("title"), body.get("paragraphs"))
        if err:
            return err
        _log(f"office/word/create path={path}")
        return jsonify({"status": "ok", **result})

    @app.route("/office/word/read", methods=["GET", "POST"])
    @require_auth
    def route_office_word_read():
        body = _json_body()
        if not isinstance(body, dict):
            body = {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        if not os.path.isfile(path):
            return jsonify({"error": "file not found", "path": path}), 404
        result, err = _word_read(path)
        if err:
            return err
        return jsonify(result)

    # ----------------------- PowerPoint file engine ----------------------
    @app.route("/office/ppt/create", methods=["POST"])
    @require_auth
    def route_office_ppt_create():
        body = _json_body()
        if not isinstance(body, dict):
            body = {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        result, err = _ppt_create(path, body.get("slides"))
        if err:
            return err
        _log(f"office/ppt/create path={path}")
        return jsonify({"status": "ok", **result})

    @app.route("/office/ppt/read", methods=["GET", "POST"])
    @require_auth
    def route_office_ppt_read():
        body = _json_body()
        if not isinstance(body, dict):
            body = {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        if not os.path.isfile(path):
            return jsonify({"error": "file not found", "path": path}), 404
        result, err = _ppt_read(path)
        if err:
            return err
        return jsonify(result)

    # ----------------------- Excel COM engine (live) ---------------------
    def _com_excel():
        if os.name != "nt":
            return None, (jsonify({"error": "Windows-only endpoint"}), 501)
        if not _HAS_COM:
            return None, _missing_engine("Excel COM", "pywin32")
        return win32com.client, None

    @app.route("/office/excel/com/open", methods=["POST"])
    @require_auth
    def route_office_excel_com_open():
        com, err = _com_excel()
        if err:
            return err
        body = _json_body() or {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        try:
            excel = com.Dispatch("Excel.Application")
            excel.Visible = bool(body.get("visible", True))
            wb = excel.Workbooks.Open(path)
            return jsonify({"status": "ok", "path": path, "sheet": wb.ActiveSheet.Name})
        except Exception as exc:
            return jsonify({"error": str(exc)}), 500

    @app.route("/office/excel/com/set", methods=["POST"])
    @require_auth
    def route_office_excel_com_set():
        com, err = _com_excel()
        if err:
            return err
        body = _json_body() or {}
        path = _resolve_path(body.get("path"))
        if not path:
            return _missing_field("path")
        cell = body.get("cell")
        value = body.get("value")
        if cell is None:
            return _missing_field("cell")
        sheet = body.get("sheet")
        try:
            excel = com.Dispatch("Excel.Application")
            excel.Visible = bool(body.get("visible", False))
            wb = excel.Workbooks.Open(path)
            ws = wb.Worksheets(sheet) if sheet else wb.ActiveSheet
            ws.Range(cell).Value = value
            if bool(body.get("save", True)):
                wb.Save()
            if not bool(body.get("keep_open", False)):
                wb.Close(SaveChanges=False)
                excel.Quit()
            return jsonify({"status": "ok", "path": path, "cell": cell})
        except Exception as exc:
            return jsonify({"error": str(exc)}), 500

    @app.route("/office/excel/com/close", methods=["POST"])
    @require_auth
    def route_office_excel_com_close():
        com, err = _com_excel()
        if err:
            return err
        try:
            excel = com.Dispatch("Excel.Application")
            if excel.Workbooks.Count > 0:
                excel.ActiveWorkbook.Close(SaveChanges=False)
            excel.Quit()
            return jsonify({"status": "ok"})
        except Exception as exc:
            return jsonify({"error": str(exc)}), 500
